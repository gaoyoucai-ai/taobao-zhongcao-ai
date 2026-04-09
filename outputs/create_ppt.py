from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 比例
prs.slide_height = Inches(7.5)

# 定义颜色
NAVY = RgbColor(28, 40, 51)      # #1C2833
SLATE = RgbColor(46, 64, 83)     # #2E4053
TEAL = RgbColor(94, 168, 167)    # #5EA8A7
DARK_TEAL = RgbColor(39, 120, 132)  # #277884
CORAL = RgbColor(254, 68, 71)    # #FE4447
RED = RgbColor(192, 57, 43)      # #C0392B
ORANGE = RgbColor(243, 156, 18)  # #F39C12
YELLOW = RgbColor(241, 196, 15)  # #F1C40F
WINE = RgbColor(93, 29, 46)      # #5D1D2E
CRIMSON = RgbColor(149, 18, 51)  # #951233
GRAY_BG = RgbColor(244, 246, 246)  # #F4F6F6

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式
    
    # 背景色块
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), Inches(0), 
        prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(255, 255, 255)
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(11.333), Inches(1))
    tf = info_box.text_frame
    p = tf.paragraphs[0]
    p.text = "汇报对象：P9 管理层    |    汇报时间：2026 年 4 月"
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(170, 183, 184)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_metrics_slide(prs):
    """添加核心成果页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心成果：超额完成目标 233%"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # 三个指标卡片
    cards = [
        {"title": "目标 CPUV", "value": "20 元", "subtitle": "", "color": TEAL},
        {"title": "实际 CPUV", "value": "8.57 元", "subtitle": "▼ 57% 优于目标", "color": CORAL},
        {"title": "投入产出比", "value": "233%", "subtitle": "超额完成", "color": ORANGE},
    ]
    
    for i, card in enumerate(cards):
        x = Inches(0.5 + i * 4.3)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.5), Inches(4), Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = card["color"]
        shape.line.fill.background()
        
        # 卡片文字
        tf = shape.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = card["title"]
        p1.font.size = Pt(16)
        p1.font.color.rgb = RgbColor(255, 255, 255)
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = card["value"]
        p2.font.size = Pt(54)
        p2.font.bold = True
        p2.font.color.rgb = RgbColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)
        
        if card["subtitle"]:
            p3 = tf.add_paragraph()
            p3.text = card["subtitle"]
            p3.font.size = Pt(20)
            p3.font.color.rgb = RgbColor(255, 255, 255)
            p3.alignment = PP_ALIGN.CENTER
            p3.space_before = Pt(10)
    
    # 关键发现框
    discovery_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.8), Inches(12.333), Inches(2)
    )
    discovery_box.fill.solid()
    discovery_box.fill.fore_color.rgb = GRAY_BG
    discovery_box.line.fill.background()
    
    # 左侧彩条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.8), Inches(0.3), Inches(2)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_TEAL
    bar.line.fill.background()
    
    tf = discovery_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    
    p = tf.paragraphs[0]
    p.text = "关键发现：照明灯具类目在 B 站存在显著流量红利，数字科技类 UP 主意外成为最优解"
    p.font.size = Pt(18)
    p.font.color.rgb = SLATE
    
    return slide

def add_case_study_slide(prs):
    """添加标杆案例页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "标杆案例：数玩工场 CPUV 6.54 元"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # 左侧达人画像
    profile_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(5.8), Inches(2.8)
    )
    profile_box.fill.solid()
    profile_box.fill.fore_color.rgb = GRAY_BG
    profile_box.line.fill.background()
    
    tf = profile_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "达人画像"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SLATE
    
    profile_data = [
        "粉丝数：10.3 万（科技数码垂类）",
        "粉丝画像：25-40 岁男性，房主/准房主",
        "内容风格：参数导向、实测对比"
    ]
    for item in profile_data:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(44, 62, 80)
        p.space_before = Pt(8)
        p.level = 0
    
    # 左侧投放效果
    results_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.8)
    )
    results_box.fill.solid()
    results_box.fill.fore_color.rgb = GRAY_BG
    results_box.line.fill.background()
    
    tf = results_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "投放效果"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = SLATE
    
    results_data = [
        "CPUV：6.54 元（行业平均 1/3）",
        "进店率：1.84%（行业平均 3 倍）",
        "转化数：16 单"
    ]
    for item in results_data:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(44, 62, 80)
        p.space_before = Pt(8)
        p.level = 0
    
    # 右侧成功公式
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(1.5), Inches(6.333), Inches(5.8)
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = WINE
    formula_box.line.fill.background()
    
    tf = formula_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "成功公式拆解"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    formula_items = [
        ("✓ 人设完美匹配", "科技粉 = 家装决策者（男性 homeowner）"),
        ("✓ 参数完整传达", "Ra>90、RG0 无频闪、智能控制等 4+ 核心指标"),
        ("✓ 真实场景展示", "安装过程实拍，非棚拍广告"),
        ("✓ 理性评论区运营", "回复参数问题，做竞品对比")
    ]
    
    for title, desc in formula_items:
        p = tf.add_paragraph()
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)
        p.space_before = Pt(12)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(13)
        p_desc.font.color.rgb = RgbColor(240, 240, 240)
        p_desc.level = 1
    
    # 底部洞察框
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.5), Inches(12.333), Inches(1.3)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = RgbColor(255, 243, 205)
    insight_box.line.fill.background()
    
    # 左侧黄条
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(7.5), Inches(0.25), Inches(1.3)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()
    
    tf = insight_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "洞察：不是"种草"逻辑，而是"求知"逻辑——B 站用户要的是购买依据，不是情感共鸣"
    p.font.size = Pt(15)
    p.font.color.rgb = RgbColor(133, 100, 4)
    p.font.bold = True
    
    return slide

def add_methodology_slide(prs):
    """添加方法论页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "可复制的方法论：三大核心要素"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # 三列要素
    elements = [
        {
            "num": "01",
            "title": "人选对 > 流量大",
            "items": ["垂直领域腰部达人（5-50 万粉）", "粉丝画像与家装决策者重合", "历史内容有参数讨论氛围"],
            "color": TEAL,
            "title_color": DARK_TEAL
        },
        {
            "num": "02",
            "title": "参数全 > 氛围好",
            "items": ["至少 3 个核心技术指标", "必须有对比环节（新旧/开关）", "优缺点并陈增加可信度"],
            "color": CORAL,
            "title_color": RED
        },
        {
            "num": "03",
            "title": "运营深 > 发布完",
            "items": ["评论区置顶参数表 + 购买链接", "2 小时内回复前 10 条评论", "主动发起参数讨论话题"],
            "color": ORANGE,
            "title_color": YELLOW
        }
    ]
    
    for i, elem in enumerate(elements):
        x = Inches(0.5 + i * 4.35)
        
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, Inches(1.5), Inches(4.1), Inches(3.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = RgbColor(255, 255, 255)
        card.line.width = Pt(3)
        card.line.fore_color.rgb = elem["color"]
        
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        
        # 编号
        p_num = tf.paragraphs[0]
        p_num.text = elem["num"]
        p_num.font.size = Pt(48)
        p_num.font.bold = True
        p_num.font.color.rgb = elem["color"]
        
        # 标题
        p_title = tf.add_paragraph()
        p_title.text = elem["title"]
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = elem["title_color"]
        p_title.space_before = Pt(10)
        
        # 列表项
        for item in elem["items"]:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = RgbColor(44, 62, 80)
            p.space_before = Pt(6)
            p.level = 0
    
    # 底部公式框
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.5), Inches(5.2), Inches(6.333), Inches(2.1)
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = NAVY
    formula_box.line.fill.background()
    
    tf = formula_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "高转化内容公式"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    formula_text = "专业人设 × 参数完整 × 真实场景 × 对比环节 × 购买引导 × 评论运营"
    p_formula = tf.add_paragraph()
    p_formula.text = formula_text
    p_formula.font.size = Pt(14)
    p_formula.font.color.rgb = RgbColor(255, 255, 255)
    p_formula.space_before = Pt(10)
    p_formula.alignment = PP_ALIGN.CENTER
    
    return slide

def add_action_plan_slide(prs):
    """添加行动计划页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "规模化复制：下一步行动计划"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # 立即执行
    immediate_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(6), Inches(2.2)
    )
    immediate_box.fill.solid()
    immediate_box.fill.fore_color.rgb = GRAY_BG
    immediate_box.line.fill.background()
    
    # 左侧绿条
    bar1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(0.25), Inches(2.2)
    )
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = TEAL
    bar1.line.fill.background()
    
    tf = immediate_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "立即执行（本周）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEAL
    
    immediate_items = [
        "✅ 联系数玩工场、会发光的微微二次合作",
        "✅ 暂停圆圈脸脸合作（0 转化排查）",
        "✅ 优化 Brief 模板（加入参数传达硬要求）"
    ]
    for item in immediate_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(44, 62, 80)
        p.space_before = Pt(6)
        p.level = 0
    
    # 中期建设
    midterm_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.833), Inches(1.5), Inches(6), Inches(2.2)
    )
    midterm_box.fill.solid()
    midterm_box.fill.fore_color.rgb = GRAY_BG
    midterm_box.line.fill.background()
    
    # 左侧橙条
    bar2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.833), Inches(1.5), Inches(0.25), Inches(2.2)
    )
    bar2.fill.solid()
    bar2.fill.fore_color.rgb = ORANGE
    bar2.line.fill.background()
    
    tf = midterm_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "中期建设（1 个月）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = YELLOW
    
    midterm_items = [
        "📋 建立 20-30 位优质达人 Whitelist",
        "📚 制作内容培训材料（产品手册 + 脚本模板）",
        "📊 搭建数据监控体系（Day3/Day7 预警线）"
    ]
    for item in midterm_items:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(44, 62, 80)
        p.space_before = Pt(6)
        p.level = 0
    
    # 长期战略
    longterm_box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4), Inches(12.333), Inches(2.3)
    )
    longterm_box.fill.solid()
    longterm_box.fill.fore_color.rgb = GRAY_BG
    longterm_box.line.fill.background()
    
    # 左侧红条
    bar3 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4), Inches(0.25), Inches(2.3)
    )
    bar3.fill.solid()
    bar3.fill.fore_color.rgb = CRIMSON
    bar3.line.fill.background()
    
    tf = longterm_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "长期战略（3 个月）"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WINE
    
    # 三个子项目
    sub_items = [
        ("达人顾问团", "邀请头部达人参与产品前期策划"),
        ("内容资产库", "积累 10-20 个标杆案例库"),
        ("形式创新", "直播测评、系列短剧、跨界联名")
    ]
    
    for i, (title, desc) in enumerate(sub_items):
        p = tf.add_paragraph()
        p.text = f"  • {title}：{desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(44, 62, 80)
        p.space_before = Pt(8)
    
    # 底部预期收益
    benefit_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.5), Inches(12.333), Inches(1.3)
    )
    benefit_box.fill.solid()
    benefit_box.fill.fore_color.rgb = TEAL
    benefit_box.line.fill.background()
    
    tf = benefit_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)
    
    p = tf.paragraphs[0]
    p.text = "预期收益"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    
    p_benefit = tf.add_paragraph()
    p_benefit.text = "规模化复制后，预计整体 CPUV 可稳定在 8-12 元区间，相比行业平均 25-30 元，节省营销预算 60%+"
    p_benefit.font.size = Pt(15)
    p_benefit.font.color.rgb = RgbColor(255, 255, 255)
    p_benefit.space_before = Pt(8)
    
    return slide

# 创建所有幻灯片
print("Creating presentation...")

# Slide 1: 封面
add_title_slide(prs, "雷士照明 KOL 投放复盘", "从个案成功到可复制的方法论")
print("✓ Slide 1: Cover")

# Slide 2: 核心成果
add_metrics_slide(prs)
print("✓ Slide 2: Core Metrics")

# Slide 3: 标杆案例
add_case_study_slide(prs)
print("✓ Slide 3: Case Study")

# Slide 4: 方法论
add_methodology_slide(prs)
print("✓ Slide 4: Methodology")

# Slide 5: 行动计划
add_action_plan_slide(prs)
print("✓ Slide 5: Action Plan")

# 保存文件
output_path = r"C:\Users\高有才\.qoderwork\workspace\mncv91pk7f82ff12\outputs\雷士照明 KOL 投放复盘-P9 汇报.pptx"
prs.save(output_path)
print(f"\n✓✓✓ Presentation saved to: {output_path}")
